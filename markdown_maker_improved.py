```python
#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter using MarkItDown

A comprehensive tool for converting PowerPoint presentations to Markdown format
with intelligent handling of complex slides through image generation and AI interpretation.

This production-ready implementation includes:
- Secure COM object handling with proper resource cleanup
- Comprehensive input validation and sanitization
- Async/await pattern for I/O operations
- Type safety with full annotations
- Configurable logging and error handling
- Path traversal protection
- Memory optimization for large presentations
"""

import argparse
import asyncio
import json
import logging
import os
import sys
import tempfile
import time
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any, AsyncGenerator
import uuid
import yaml
from concurrent.futures import ThreadPoolExecutor

# Third-party imports with error handling
try:
    from PIL import Image, ImageDraw
    import comtypes.client
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.picture import Picture
    from pptx.shapes.table import Table
    from pptx.shapes.chart import Chart
    from markitdown import MarkItDown
    import aiofiles
except ImportError as e:
    print(f"Missing required dependency: {e}")
    print("Please install requirements: pip install -r requirements.txt")
    sys.exit(1)


def setup_logging(verbose: bool = False, log_file: Optional[Path] = None) -> logging.Logger:
    """Configure logging with proper formatting and handlers."""
    level = logging.DEBUG if verbose else logging.INFO
    
    # Create formatter
    formatter = logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    
    # Configure root logger
    root_logger = logging.getLogger()
    root_logger.setLevel(level)
    
    # Clear existing handlers
    for handler in root_logger.handlers[:]:
        root_logger.removeHandler(handler)
    
    # Console handler
    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setFormatter(formatter)
    root_logger.addHandler(console_handler)
    
    # File handler if specified
    if log_file:
        log_file.parent.mkdir(parents=True, exist_ok=True)
        file_handler = logging.FileHandler(log_file)
        file_handler.setFormatter(formatter)
        root_logger.addHandler(file_handler)
    
    return logging.getLogger(__name__)


@dataclass
class SlideAnalysis:
    """Data class for slide analysis results."""
    slide_number: int
    title: str
    object_count: int
    has_images: bool
    has_tables: bool
    has_charts: bool
    is_complex: bool
    shapes: List[str] = field(default_factory=list)
    text_content: str = ""
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            'slide_number': self.slide_number,
            'title': self.title,
            'object_count': self.object_count,
            'has_images': self.has_images,
            'has_tables': self.has_tables,
            'has_charts': self.has_charts,
            'is_complex': self.is_complex,
            'shapes': self.shapes,
            'text_content': self.text_content
        }


@dataclass
class ConversionConfig:
    """Configuration for PowerPoint to Markdown conversion with validation."""
    object_threshold: int = 5
    image_dpi: int = 300
    output_dir: Path = Path("output")
    generate_mermaid: bool = True
    preserve_structure: bool = True
    batch_mode: bool = False
    verbose: bool = False
    max_file_size: int = 100 * 1024 * 1024  # 100MB
    timeout_seconds: int = 300  # 5 minutes
    retry_attempts: int = 3
    
    def __post_init__(self) -> None:
        """Validate configuration parameters."""
        if self.object_threshold < 1:
            raise ValueError("object_threshold must be positive")
        
        if self.image_dpi < 72:
            raise ValueError("image_dpi must be at least 72")
        
        if self.max_file_size < 1024:  # 1KB minimum
            raise ValueError("max_file_size must be at least 1KB")
        
        if self.timeout_seconds < 1:
            raise ValueError("timeout_seconds must be positive")
        
        if self.retry_attempts < 1:
            raise ValueError("retry_attempts must be positive")
        
        # Ensure output directory is absolute and create if needed
        self.output_dir = Path(self.output_dir).resolve()
    
    @classmethod
    def from_file(cls, config_path: Path) -> 'ConversionConfig':
        """Load configuration from YAML or JSON file."""
        if not config_path.exists():
            raise FileNotFoundError(f"Configuration file not found: {config_path}")
        
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                if config_path.suffix.lower() == '.yaml' or config_path.suffix.lower() == '.yml':
                    data = yaml.safe_load(f)
                else:
                    data = json.load(f)
            
            # Convert output_dir string to Path if present
            if 'output_dir' in data:
                data['output_dir'] = Path(data['output_dir'])
            
            return cls(**data)
        except Exception as e:
            raise ValueError(f"Failed to load configuration file: {e}")


class SecurityValidator:
    """Security validation utilities."""
    
    ALLOWED_EXTENSIONS = {'.pptx', '.ppt'}
    MAX_PATH_DEPTH = 10
    
    @staticmethod
    def validate_file_path(file_path: Path, base_dir: Optional[Path] = None) -> Path:
        """Validate file path for security issues."""
        try:
            # Resolve to absolute path
            resolved_path = file_path.resolve()
            
            # Check if file exists
            if not resolved_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Check file extension
            if resolved_path.suffix.lower() not in SecurityValidator.ALLOWED_EXTENSIONS:
                raise ValueError(f"Invalid file type: {resolved_path.suffix}")
            
            # Prevent path traversal
            if base_dir:
                base_resolved = base_dir.resolve()
                try:
                    resolved_path.relative_to(base_resolved)
                except ValueError:
                    raise ValueError(f"Path outside allowed directory: {resolved_path}")
            
            # Check path depth
            if len(resolved_path.parts) > SecurityValidator.MAX_PATH_DEPTH:
                raise ValueError("Path too deep")
            
            return resolved_path
            
        except OSError as e:
            raise ValueError(f"Invalid file path: {e}")
    
    @staticmethod
    def validate_file_size(file_path: Path, max_size: int) -> None:
        """Validate file size."""
        try:
            size = file_path.stat().st_size
            if size > max_size:
                raise ValueError(f"File too large: {size} bytes (max: {max_size})")
        except OSError as e:
            raise ValueError(f"Cannot access file: {e}")
    
    @staticmethod
    def sanitize_filename(filename: str) -> str:
        """Sanitize filename for safe filesystem operations."""
        # Remove or replace dangerous characters
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        
        # Limit length
        filename = filename[:255]
        
        # Ensure not empty
        if not filename.strip():
            filename = f"unnamed_{uuid.uuid4().hex[:8]}"
        
        return filename


class PowerPointAnalyzer:
    """Analyzes PowerPoint presentations to identify slide complexity and structure."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
    
    async def analyze_presentation(self, pptx_path: Path) -> List[SlideAnalysis]:
        """
        Analyze a PowerPoint presentation and return slide analysis data.
        
        Args:
            pptx_path: Path to the PowerPoint file
            
        Returns:
            List of SlideAnalysis objects for each slide
            
        Raises:
            FileNotFoundError: If the PowerPoint file doesn't exist
            ValueError: If the file is not a valid PowerPoint presentation
        """
        # Security validation
        validated_path = SecurityValidator.validate_file_path(pptx_path)
        SecurityValidator.validate_file_size(validated_path, self.config.max_file_size)
        
        try:
            # Use thread pool for CPU-intensive analysis
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor(max_workers=1) as executor:
                analyses = await loop.run_in_executor(
                    executor, 
                    self._analyze_presentation_sync, 
                    validated_path
                )
            
            return analyses
            
        except Exception as e:
            self.logger.error(f"Failed to analyze PowerPoint presentation: {e}")
            raise ValueError(f"Failed to analyze PowerPoint presentation: {e}")
    
    def _analyze_presentation_sync(self, pptx_path: Path) -> List[SlideAnalysis]:
        """Synchronous presentation analysis (runs in thread pool)."""
        try:
            presentation = Presentation(str(pptx_path))
            analyses = []
            
            self.logger.info(f"Analyzing presentation with {len(presentation.slides)} slides")
            
            for idx, slide in enumerate(presentation.slides, 1):
                analysis = self._analyze_slide(slide, idx)
                analyses.append(analysis)
                
                if self.config.verbose:
                    self.logger.debug(f"Slide {idx}: {analysis.object_count} objects, "
                                    f"Complex: {analysis.is_complex}")
            
            return analyses
            
        except Exception as e:
            raise ValueError(f"Failed to open PowerPoint presentation: {e}")
    
    def _analyze_slide(self, slide, slide_number: int) -> SlideAnalysis:
        """Analyze a single slide for complexity and content."""
        shapes = []
        object_count = 0
        has_images = False
        has_tables = False
        has_charts = False
        title = ""
        text_content = []
        
        # Extract title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        
        # Count and categorize shapes
        for shape in slide.shapes:
            object_count += 1
            shape_type = self._get_shape_type(shape)
            shapes.append(shape_type)
            
            # Extract text content
            if hasattr(shape, 'text') and shape.text.strip():
                text_content.append(shape.text.strip())
            
            # Categorize shape types
            if isinstance(shape, Picture):
                has_images = True
            elif isinstance(shape, Table):
                has_tables = True
            elif isinstance(shape, Chart):
                has_charts = True
        
        is_complex = object_count >= self.config.object_threshold
        
        return SlideAnalysis(
            slide_number=slide_number,
            title=title or f"Slide {slide_number}",
            object_count=object_count,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
            is_complex=is_complex,
            shapes=shapes,
            text_content="\n".join(text_content)
        )
    
    def _get_shape_type(self, shape: BaseShape) -> str:
        """Get a human-readable shape type."""
        shape_type_map = {
            1: "AutoShape",
            13: "Picture", 
            19: "Table",
            3: "Chart",
            17: "TextBox",
            5: "Freeform",
            9: "Group"
        }
        return shape_type_map.get(shape.shape_type, f"Unknown({shape.shape_type})")


class SlideImageGenerator:
    """Generates high-quality images from complex PowerPoint slides."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
        self._temp_files: List[Path] = []
    
    async def generate_slide_images(
        self, 
        pptx_path: Path, 
        analyses: List[SlideAnalysis]
    ) -> Dict[int, Path]:
        """
        Generate PNG images for complex slides.
        
        Args:
            pptx_path: Path to the PowerPoint file
            analyses: List of slide analyses
            
        Returns:
            Dictionary mapping slide numbers to generated image paths
        """
        complex_slides = [a for a in analyses if a.is_complex]
        if not complex_slides:
            self.logger.info("No complex slides found, skipping image generation")
            return {}
        
        self.logger.info(f"Generating images for {len(complex_slides)} complex slides")
        
        # Ensure output directory exists
        images_dir = self.config.output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            return await self._convert_slides_to_images(pptx_path, complex_slides, images_dir)
        except Exception as e:
            self.logger.error(f"Failed to generate slide images: {e}")
            # Clean up temporary files
            await self._cleanup_temp_files()
            raise
    
    @asynccontextmanager
    async def _com_powerpoint(self, pptx_path: Path):
        """Secure COM PowerPoint context manager."""
        ppt_app = None
        presentation = None
        
        try:
            # Initialize PowerPoint application via COM
            ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
            ppt_app.Visible = False
            ppt_app.DisplayAlerts = False
            
            # Open presentation with error handling
            presentation = ppt_app.Presentations.Open(
                str(pptx_path.absolute()), 
                ReadOnly=True, 
                Untitled=True, 
                WithWindow=False
            )
            
            yield ppt_app, presentation
            
        except comtypes.COMError as e:
            self.logger.error(f"COM error: {e}")
            raise
        except Exception as e:
            self.logger.error(f"PowerPoint automation error: {e}")
            raise
        finally:
            # Ensure proper cleanup
            try:
                if presentation:
                    presentation.Close()
                if ppt_app:
                    ppt_app.Quit()
            except Exception as e:
                self.logger.warning(f"Error during COM cleanup: {e}")
    
    async def _convert_slides_to_images(
        self, 
        pptx_path: Path, 
        complex_slides: List[SlideAnalysis], 
        output_dir: Path
    ) -> Dict[int, Path]:
        """Convert specified slides to PNG images using COM automation."""
        image_paths = {}
        
        try:
            async with self._com_power