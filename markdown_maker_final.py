#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter using MarkItDown
Complete, working version with essential features
"""

import argparse
import asyncio
import logging
import os
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional

# Third-party imports (with error handling)
try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from markitdown import MarkItDown
except ImportError:
    print("Error: markitdown not installed. Run: pip install markitdown")
    sys.exit(1)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class SlideAnalysis:
    """Data class for slide analysis results."""
    slide_number: int
    title: str
    object_count: int
    has_images: bool
    has_tables: bool
    has_charts: bool
    is_complex: bool
    shapes: List[str]


@dataclass
class ConversionConfig:
    """Configuration for PowerPoint to Markdown conversion."""
    object_threshold: int = 5
    output_dir: Path = Path("output")
    verbose: bool = False


class PowerPointAnalyzer:
    """Analyzes PowerPoint presentations to identify slide complexity."""

    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")

    def analyze_presentation(self, pptx_path: Path) -> List[SlideAnalysis]:
        """Analyze a PowerPoint presentation and return slide analysis data."""
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

        if pptx_path.suffix.lower() not in ['.pptx', '.ppt']:
            raise ValueError(f"Invalid file type: {pptx_path.suffix}")

        try:
            presentation = Presentation(str(pptx_path))
            analyses = []

            self.logger.info(f"Analyzing presentation with {len(presentation.slides)} slides")

            for idx, slide in enumerate(presentation.slides, 1):
                analysis = self._analyze_slide(slide, idx)
                analyses.append(analysis)

                if self.config.verbose:
                    self.logger.debug(f"Slide {idx}: {analysis.object_count} objects, "
                                    f"Complex: {analysis.is_complex}")

            return analyses

        except Exception as e:
            raise ValueError(f"Failed to analyze PowerPoint presentation: {e}")

    def _analyze_slide(self, slide, slide_number: int) -> SlideAnalysis:
        """Analyze a single slide for complexity and content."""
        shapes = []
        object_count = 0
        has_images = False
        has_tables = False
        has_charts = False
        title = ""

        # Extract title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Count and categorize shapes
        for shape in slide.shapes:
            object_count += 1
            shape_type = self._get_shape_type(shape)
            shapes.append(shape_type)

            # Check for specific types
            if hasattr(shape, 'image'):
                has_images = True
            elif hasattr(shape, 'table'):
                has_tables = True
            elif hasattr(shape, 'chart'):
                has_charts = True

        is_complex = object_count >= self.config.object_threshold

        return SlideAnalysis(
            slide_number=slide_number,
            title=title or f"Slide {slide_number}",
            object_count=object_count,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
            is_complex=is_complex,
            shapes=shapes
        )

    def _get_shape_type(self, shape: BaseShape) -> str:
        """Get a human-readable shape type."""
        return f"Shape({shape.shape_type})"


class MarkItDownConverter:
    """Converts PowerPoint files to Markdown using MarkItDown library."""

    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
        self.markitdown = MarkItDown()

    async def convert_to_markdown(self, pptx_path: Path, analyses: List[SlideAnalysis]) -> str:
        """Convert PowerPoint to Markdown using MarkItDown."""
        try:
            self.logger.info(f"Converting {pptx_path} to Markdown using MarkItDown")

            # Use MarkItDown for conversion
            result = self.markitdown.convert(str(pptx_path))
            markdown_content = result.text_content

            # Add analysis summary
            summary = self._generate_analysis_summary(analyses)
            full_content = f"{summary}\n\n{markdown_content}"

            return full_content

        except Exception as e:
            self.logger.error(f"MarkItDown conversion failed: {e}")
            raise


    def _generate_analysis_summary(self, analyses: List[SlideAnalysis]) -> str:
        """Generate a summary of the slide analysis."""
        total_slides = len(analyses)
        complex_slides = sum(1 for a in analyses if a.is_complex)

        summary = f"""# Presentation Analysis Summary

- **Total Slides**: {total_slides}
- **Complex Slides** (≥{self.config.object_threshold} objects): {complex_slides}
- **Conversion Method**: MarkItDown AI-powered conversion

## Complex Slides Overview
"""

        for analysis in analyses:
            if analysis.is_complex:
                summary += f"- Slide {analysis.slide_number}: {analysis.title} ({analysis.object_count} objects)\n"

        return summary


class OutputManager:
    """Manages output file generation and organization."""

    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")

    async def save_markdown(self, content: str, pptx_path: Path, analyses: List[SlideAnalysis]) -> Path:
        """Save markdown content to file."""
        # Ensure output directory exists
        self.config.output_dir.mkdir(parents=True, exist_ok=True)

        # Generate output filename
        output_filename = f"{pptx_path.stem}_converted.md"
        output_path = self.config.output_dir / output_filename

        try:
            # Write markdown content
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)

            self.logger.info(f"Markdown saved to: {output_path}")
            return output_path

        except Exception as e:
            self.logger.error(f"Failed to save markdown: {e}")
            raise


class MarkdownMakerApp:
    """Main application orchestrator for PowerPoint to Markdown conversion."""

    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")

        # Initialize components
        self.analyzer = PowerPointAnalyzer(config)
        self.converter = MarkItDownConverter(config)
        self.output_manager = OutputManager(config)

    async def convert_presentation(self, pptx_path: Path) -> Dict[str, any]:
        """Convert a PowerPoint presentation to Markdown."""
        try:
            self.logger.info(f"Starting conversion of {pptx_path}")

            # Step 1: Analyze presentation
            analyses = self.analyzer.analyze_presentation(pptx_path)

            # Step 2: Convert to Markdown
            markdown_content = await self.converter.convert_to_markdown(pptx_path, analyses)

            # Step 3: Save output
            output_path = await self.output_manager.save_markdown(markdown_content, pptx_path, analyses)

            # Return results
            return {
                "status": "success",
                "input_file": str(pptx_path),
                "output_file": str(output_path),
                "total_slides": len(analyses),
                "complex_slides": sum(1 for a in analyses if a.is_complex),
                "analyses": analyses
            }

        except Exception as e:
            self.logger.error(f"Conversion failed: {e}")
            return {
                "status": "error",
                "input_file": str(pptx_path),
                "error": str(e)
            }


def main():
    """Main CLI function."""
    parser = argparse.ArgumentParser(description="Convert PowerPoint to Markdown using MarkItDown")
    parser.add_argument("input_files", nargs="+", help="PowerPoint files to convert")
    parser.add_argument("--threshold", type=int, default=5, help="Object threshold for complexity (default: 5)")
    parser.add_argument("--output-dir", type=str, default="output", help="Output directory (default: output)")
    parser.add_argument("--verbose", action="store_true", help="Verbose logging")

    args = parser.parse_args()

    # Configure logging
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # Create configuration
    config = ConversionConfig(
        object_threshold=args.threshold,
        output_dir=Path(args.output_dir),
        verbose=args.verbose
    )

    # Create app instance
    app = MarkdownMakerApp(config)

    async def convert_files():
        """Convert all input files."""
        results = []

        for input_file in args.input_files:
            pptx_path = Path(input_file)

            if not pptx_path.exists():
                logger.error(f"File not found: {pptx_path}")
                continue

            result = await app.convert_presentation(pptx_path)
            results.append(result)

            if result["status"] == "success":
                logger.info(f"✅ Converted: {result['input_file']} → {result['output_file']}")
            else:
                logger.error(f"❌ Failed: {result['input_file']} - {result['error']}")

        # Summary
        successful = sum(1 for r in results if r["status"] == "success")
        total = len(results)

        print(f"\n🎉 Conversion complete: {successful}/{total} files successful")

        return results

    # Run conversion
    try:
        asyncio.run(convert_files())
    except KeyboardInterrupt:
        logger.info("Conversion interrupted by user")
    except Exception as e:
        logger.error(f"Application error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()