#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter using MarkItDown

A comprehensive tool for converting PowerPoint presentations to Markdown format
with intelligent handling of complex slides through image generation and AI interpretation.
"""

import argparse
import asyncio
import json
import logging
import os
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Any
import yaml

# Third-party imports
from PIL import Image, ImageDraw
import comtypes.client
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.table import Table
from pptx.shapes.chart import Chart
from markitdown import MarkItDown

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class SlideAnalysis:
    """Data class for slide analysis results."""
    slide_number: int
    title: str
    object_count: int
    has_images: bool
    has_tables: bool
    has_charts: bool
    is_complex: bool
    shapes: List[str] = field(default_factory=list)


@dataclass
class ConversionConfig:
    """Configuration for PowerPoint to Markdown conversion."""
    object_threshold: int = 5
    image_dpi: int = 300
    output_dir: Path = Path("output")
    generate_mermaid: bool = True
    preserve_structure: bool = True
    batch_mode: bool = False
    verbose: bool = False


class PowerPointAnalyzer:
    """Analyzes PowerPoint presentations to identify slide complexity and structure."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
    
    def analyze_presentation(self, pptx_path: Path) -> List[SlideAnalysis]:
        """
        Analyze a PowerPoint presentation and return slide analysis data.
        
        Args:
            pptx_path: Path to the PowerPoint file
            
        Returns:
            List of SlideAnalysis objects for each slide
            
        Raises:
            FileNotFoundError: If the PowerPoint file doesn't exist
            ValueError: If the file is not a valid PowerPoint presentation
        """
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        try:
            presentation = Presentation(str(pptx_path))
            analyses = []
            
            self.logger.info(f"Analyzing presentation with {len(presentation.slides)} slides")
            
            for idx, slide in enumerate(presentation.slides, 1):
                analysis = self._analyze_slide(slide, idx)
                analyses.append(analysis)
                
                if self.config.verbose:
                    self.logger.debug(f"Slide {idx}: {analysis.object_count} objects, "
                                    f"Complex: {analysis.is_complex}")
            
            return analyses
            
        except Exception as e:
            raise ValueError(f"Failed to analyze PowerPoint presentation: {e}")
    
    def _analyze_slide(self, slide, slide_number: int) -> SlideAnalysis:
        """Analyze a single slide for complexity and content."""
        shapes = []
        object_count = 0
        has_images = False
        has_tables = False
        has_charts = False
        title = ""
        
        # Extract title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        
        # Count and categorize shapes
        for shape in slide.shapes:
            object_count += 1
            shape_type = self._get_shape_type(shape)
            shapes.append(shape_type)
            
            if isinstance(shape, Picture):
                has_images = True
            elif isinstance(shape, Table):
                has_tables = True
            elif isinstance(shape, Chart):
                has_charts = True
        
        is_complex = object_count >= self.config.object_threshold
        
        return SlideAnalysis(
            slide_number=slide_number,
            title=title or f"Slide {slide_number}",
            object_count=object_count,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
            is_complex=is_complex,
            shapes=shapes
        )
    
    def _get_shape_type(self, shape: BaseShape) -> str:
        """Get a human-readable shape type."""
        shape_type_map = {
            1: "AutoShape",
            13: "Picture", 
            19: "Table",
            3: "Chart",
            17: "TextBox",
            5: "Freeform",
            9: "Group"
        }
        return shape_type_map.get(shape.shape_type, f"Unknown({shape.shape_type})")


class SlideImageGenerator:
    """Generates high-quality images from complex PowerPoint slides."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
    
    async def generate_slide_images(
        self, 
        pptx_path: Path, 
        analyses: List[SlideAnalysis]
    ) -> Dict[int, Path]:
        """
        Generate PNG images for complex slides.
        
        Args:
            pptx_path: Path to the PowerPoint file
            analyses: List of slide analyses
            
        Returns:
            Dictionary mapping slide numbers to generated image paths
        """
        complex_slides = [a for a in analyses if a.is_complex]
        if not complex_slides:
            self.logger.info("No complex slides found, skipping image generation")
            return {}
        
        self.logger.info(f"Generating images for {len(complex_slides)} complex slides")
        
        # Ensure output directory exists
        images_dir = self.config.output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            return await self._convert_slides_to_images(pptx_path, complex_slides, images_dir)
        except Exception as e:
            self.logger.error(f"Failed to generate slide images: {e}")
            raise
    
    async def _convert_slides_to_images(
        self, 
        pptx_path: Path, 
        complex_slides: List[SlideAnalysis], 
        output_dir: Path
    ) -> Dict[int, Path]:
        """Convert specified slides to PNG images using COM automation."""
        image_paths = {}
        
        try:
            # Initialize PowerPoint application via COM
            ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
            ppt_app.Visible = False
            
            # Open presentation
            presentation = ppt_app.Presentations.Open(str(pptx_path.absolute()))
            
            for slide_analysis in complex_slides:
                slide_num = slide_analysis.slide_number
                image_filename = f"slide_{slide_num:03d}.png"
                image_path = output_dir / image_filename
                
                try:
                    # Export slide as PNG
                    slide = presentation.Slides(slide_num)
                    slide.Export(str(image_path), "PNG", 1920, 1080)  # HD resolution
                    
                    # Optimize the generated image
                    await self._optimize_image(image_path)
                    
                    image_paths[slide_num] = image_path
                    self.logger.debug(f"Generated image for slide {slide_num}: {image_path}")
                    
                except Exception as e:
                    self.logger.warning(f"Failed to export slide {slide_num}: {e}")
            
            # Cleanup
            presentation.Close()
            ppt_app.Quit()
            
        except Exception as e:
            self.logger.error(f"COM automation failed: {e}")
            # Fallback to alternative method if available
            
        return image_paths
    
    async def _optimize_image(self, image_path: Path) -> None:
        """Optimize the generated image for better quality and size."""
        try:
            with Image.open(image_path) as img:
                # Ensure RGB mode
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Apply optimization
                img.save(image_path, "PNG", optimize=True, quality=95)
                
        except Exception as e:
            self.logger.warning(f"Failed to optimize image {image_path}: {e}")


class MarkItDownConverter:
    """Integrates with MarkItDown library for AI-powered conversion."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.markitdown = MarkItDown()
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
    
    async def convert_to_markdown(
        self, 
        pptx_path: Path, 
        analyses: List[SlideAnalysis],
        image_paths: Dict[int, Path]
    ) -> str:
        """
        Convert PowerPoint presentation to Markdown format.
        
        Args:
            pptx_path: Path to the PowerPoint file
            analyses: List of slide analyses
            image_paths: Dictionary of slide numbers to image paths
            
        Returns:
            Generated Markdown content
        """
        self.logger.info("Converting presentation to Markdown")
        
        markdown_content = []
        
        # Add title
        presentation_name = pptx_path.stem
        markdown_content.append(f"# {presentation_name}\n")
        
        # Convert slides
        for analysis in analyses:
            slide_md = await self._convert_slide(analysis, image_paths)
            markdown_content.append(slide_md)
        
        return "\n\n".join(markdown_content)
    
    async def _convert_slide(
        self, 
        analysis: SlideAnalysis, 
        image_paths: Dict[int, Path]
    ) -> str:
        """Convert a single slide to Markdown."""
        slide_md = [f"## {analysis.title}"]
        
        if analysis.is_complex and analysis.slide_number in image_paths:
            # Use image for complex slides
            image_path = image_paths[analysis.slide_number]
            relative_path = Path("images") / image_path.name
            
            slide_md.append(f"![{analysis.title}]({relative_path})")
            
            # Add AI interpretation if available
            try:
                ai_description = await self._get_ai_interpretation(image_path)
                if ai_description:
                    slide_md.append(ai_description)
            except Exception as e:
                self.logger.warning(f"Failed to get AI interpretation for slide {analysis.slide_number}: {e}")
        
        else:
            # Add slide metadata for simple slides
            slide_md.append(f"*Slide {analysis.slide_number} - {analysis.object_count} objects*")
            
            if analysis.has_charts and self.config.generate_mermaid:
                slide_md.append(self._generate_mermaid_placeholder(analysis))
        
        return "\n\n".join(slide_md)
    
    async def _get_ai_interpretation(self, image_path: Path) -> Optional[str]:
        """Get AI interpretation of slide image using MarkItDown."""
        try:
            # Use MarkItDown to interpret the image
            result = self.markitdown.convert(str(image_path))
            return result.text_content if result else None
        except Exception as e:
            self.logger.debug(f"AI interpretation failed for {image_path}: {e}")
            return None
    
    def _generate_mermaid_placeholder(self, analysis: SlideAnalysis) -> str:
        """Generate Mermaid diagram placeholder for chart-containing slides."""
        return f"""```mermaid
graph TD
    A[{analysis.title}] --> B[Chart Data]
    B --> C[Analysis Results]
```

*Note: Mermaid diagram placeholder for slide {analysis.slide_number}*"""


class OutputManager:
    """Manages output file organization and Markdown generation."""
    
    def __init__(self, config: ConversionConfig):
        self.config = config
        self.logger = logging.getLogger(f"{__name__}.{self.__class__.__name__}")
    
    async def save_markdown(
        self, 
        content: str, 
        pptx_path: Path,
        analyses: List[SlideAnalysis]
    ) -> Path:
        """
        Save Markdown content to file with proper organization.
        
        Args:
            content: Generated Markdown content
            pptx_path: Original PowerPoint file path
            analyses: List of slide analyses
            
        Returns:
            Path to the saved Markdown file
        """
        # Create output directory
        self.config.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Generate output filename
        output_filename = f"{pptx_path.stem}.md"
        output_path = self.config.output_dir / output_filename
        
        # Add metadata header
        metadata = self._generate_metadata(pptx_path, analyses)
        full_content = metadata + "\n\n" + content
        
        # Write to file
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_content)
            
            self.logger.info(f"Markdown saved to: {output_path}")
            return output_path
            
        except Exception as e:
            raise IOError(f"Failed to save Markdown file: {e}")
    
    def _generate_metadata(self, pptx_path: Path, analyses: List[SlideAnalysis]) -> str:
        """Generate metadata header for the Markdown file."""
        complex_count = sum(1 for a in analyses if a.is_complex)
        
        metadata = f"""---
source_file: {pptx_path.name}
total_slides: {len(analyses)}
complex_slides: {complex_count}
conversion_date: {asyncio.get_event_loop().time()}
generator: MarkdownMaker v1.0
---"""
        
        return metadata
    
    def generate_summary_report(self, analyses: List[SlideAnalysis]) -> str:
        """Generate a summary report of the conversion."""
        total_slides = len(analyses)
        complex_slides = sum(1 for a in analyses if a.is_complex)
        slides_with_images = sum(1 for a in analyses if a.has_images)
        slides_with_tables = sum(1 for a in analyses if a.has_tables)
        slides_with_charts = sum(1 for a in analyses if a.has_charts)
        
        report = f"""
# Conversion Summary

- **Total Slides**: {total_slides}
- **Complex Slides** (≥{self.config.object_threshold} objects): {complex_slides}
- **Slides with Images**: {slides_with_images}
- **Slides with Tables**: {slides_with_tables}
- **Slides with Charts**: {slides_with_charts}

## Complex Slides Details

| Slide | Title | Objects | Images | Tables | Charts |
|-------|-------|---------|---------|---------|---------|
"""
        
        for analysis in analyses:
            if analysis.is_complex:
                report += f"| {analysis.slide_number} | {analysis.title[:30]}... | {analysis.object_count