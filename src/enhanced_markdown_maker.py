#!/usr/bin/env python3
"""
Enhanced MarkdownMaker - PowerPoint to Markdown conversion system
Built on Microsoft MarkItDown library with enhanced object counting and Mermaid conversion
"""

import os
import json
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, asdict
from datetime import datetime

try:
    from markitdown import MarkItDown
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw
    import io
except ImportError as e:
    raise ImportError(
        f"Missing required dependency: {e}. "
        "Please install: pip install markitdown python-pptx Pillow"
    )


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@dataclass
class SlideObject:
    """Represents an object found in a PowerPoint slide"""
    object_type: str
    description: str
    position: Optional[Dict[str, float]] = None
    size: Optional[Dict[str, float]] = None


@dataclass
class SlideAnalysis:
    """Analysis results for a single slide"""
    slide_number: int
    object_count: int
    objects: List[SlideObject]
    has_image: bool
    has_chart: bool
    has_table: bool
    complexity_score: float
    should_save_as_image: bool
    title: Optional[str] = None


@dataclass
class ConversionMetadata:
    """Metadata for the entire conversion process"""
    source_file: str
    total_slides: int
    slides_as_images: int
    slides_as_text: int
    mermaid_diagrams: int
    conversion_time: float
    timestamp: str
    slide_analyses: List[SlideAnalysis]


class EnhancedMarkdownMaker:
    """
    Main class for enhanced PowerPoint to Markdown conversion.
    Uses MarkItDown library with custom object counting and Mermaid generation.
    """

    def __init__(
        self,
        llm_client: Optional[Any] = None,
        output_dir: str = "output",
        object_threshold: int = 5,
        llm_model: str = "gpt-4o",
        enable_mermaid: bool = True
    ):
        """
        Initialize the Enhanced MarkdownMaker.

        Args:
            llm_client: OpenAI-compatible LLM client for image analysis
            output_dir: Directory for output files
            object_threshold: Number of objects to trigger image conversion
            llm_model: LLM model to use for analysis
            enable_mermaid: Enable Mermaid diagram generation
        """
        self.output_dir = Path(output_dir)
        self.object_threshold = object_threshold
        self.llm_model = llm_model
        self.enable_mermaid = enable_mermaid

        # Create output directories
        self.slides_dir = self.output_dir / "slides"
        self.mermaid_dir = self.output_dir / "mermaid"
        self.slides_dir.mkdir(parents=True, exist_ok=True)
        self.mermaid_dir.mkdir(parents=True, exist_ok=True)

        # Initialize MarkItDown with custom prompt if LLM client provided
        if llm_client:
            self.markitdown = MarkItDown(
                llm_client=llm_client,
                llm_model=llm_model,
                llm_prompt=self._create_analysis_prompt()
            )
        else:
            self.markitdown = MarkItDown()
            logger.warning("No LLM client provided. Advanced analysis features disabled.")

        logger.info(f"EnhancedMarkdownMaker initialized with output_dir={output_dir}")

    def _create_analysis_prompt(self) -> str:
        """Create comprehensive LLM prompt for slide analysis"""
        return """
        Analyze this PowerPoint slide image and provide structured analysis:

        1. Count all distinct visual objects:
           - Shapes (rectangles, circles, arrows, etc.)
           - Text boxes and labels
           - Images and icons
           - Charts and graphs
           - Tables
           - Connectors and lines

        2. Identify the slide's purpose and content

        3. If this appears to be a diagram or flowchart, provide Mermaid syntax

        4. Return as JSON with this structure:
        {
            "object_count": <number>,
            "title": "<slide title if visible>",
            "description": "<detailed content description>",
            "slide_type": "<diagram|flowchart|text|mixed|chart>",
            "mermaid": "<mermaid syntax or null>",
            "objects": [
                {
                    "type": "<shape|text|image|chart|table>",
                    "description": "<what this object represents>"
                }
            ],
            "key_points": ["<bullet points of main content>"]
        }

        Be precise with object counting. Each distinct visual element counts as one object.
        """

    def convert_powerpoint(
        self,
        pptx_path: str,
        output_filename: Optional[str] = None
    ) -> Tuple[str, ConversionMetadata]:
        """
        Convert PowerPoint to Markdown with enhanced analysis.

        Args:
            pptx_path: Path to PowerPoint file
            output_filename: Optional custom output filename

        Returns:
            Tuple of (markdown_content, metadata)
        """
        start_time = datetime.now()
        logger.info(f"Starting conversion of {pptx_path}")

        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

        # Analyze presentation structure
        prs = Presentation(str(pptx_path))
        slide_analyses = self._analyze_all_slides(prs)

        # Generate markdown content
        markdown_content = self._generate_enhanced_markdown(
            prs, pptx_path, slide_analyses
        )

        # Create metadata
        end_time = datetime.now()
        metadata = ConversionMetadata(
            source_file=str(pptx_path),
            total_slides=len(prs.slides),
            slides_as_images=sum(1 for a in slide_analyses if a.should_save_as_image),
            slides_as_text=sum(1 for a in slide_analyses if not a.should_save_as_image),
            mermaid_diagrams=0,  # Will be updated during Mermaid generation
            conversion_time=(end_time - start_time).total_seconds(),
            timestamp=start_time.isoformat(),
            slide_analyses=slide_analyses
        )

        # Save markdown file
        if output_filename is None:
            output_filename = f"{pptx_path.stem}.md"

        output_path = self.output_dir / output_filename
        output_path.write_text(markdown_content, encoding='utf-8')
        logger.info(f"Markdown saved to {output_path}")

        # Save metadata
        metadata_path = self.output_dir / "metadata.json"
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(asdict(metadata), f, indent=2, default=str)
        logger.info(f"Metadata saved to {metadata_path}")

        return markdown_content, metadata

    def _analyze_all_slides(self, prs: Presentation) -> List[SlideAnalysis]:
        """Analyze all slides in the presentation"""
        analyses = []

        for idx, slide in enumerate(prs.slides, 1):
            analysis = self._analyze_slide_objects(slide, idx)
            analyses.append(analysis)
            logger.debug(f"Slide {idx}: {analysis.object_count} objects, "
                        f"complexity={analysis.complexity_score:.2f}")

        return analyses

    def _analyze_slide_objects(self, slide, slide_num: int) -> SlideAnalysis:
        """
        Analyze objects in a single slide.

        Args:
            slide: python-pptx Slide object
            slide_num: Slide number (1-indexed)

        Returns:
            SlideAnalysis object
        """
        objects = []
        has_image = False
        has_chart = False
        has_table = False
        title = None

        # Extract title if available
        if slide.shapes.title:
            title = slide.shapes.title.text

        for shape in slide.shapes:
            obj = self._classify_shape(shape)
            if obj:
                objects.append(obj)

                # Update flags
                if obj.object_type == "image":
                    has_image = True
                elif obj.object_type == "chart":
                    has_chart = True
                elif obj.object_type == "table":
                    has_table = True

        # Calculate complexity score
        complexity_score = self._calculate_complexity(
            len(objects), has_image, has_chart, has_table
        )

        # Determine if should save as image
        should_save = len(objects) >= self.object_threshold

        return SlideAnalysis(
            slide_number=slide_num,
            object_count=len(objects),
            objects=objects,
            has_image=has_image,
            has_chart=has_chart,
            has_table=has_table,
            complexity_score=complexity_score,
            should_save_as_image=should_save,
            title=title
        )

    def _classify_shape(self, shape) -> Optional[SlideObject]:
        """Classify a shape and extract relevant information"""
        try:
            shape_type = shape.shape_type

            # Get position and size
            position = {
                "left": float(shape.left) if hasattr(shape, 'left') else 0,
                "top": float(shape.top) if hasattr(shape, 'top') else 0
            }
            size = {
                "width": float(shape.width) if hasattr(shape, 'width') else 0,
                "height": float(shape.height) if hasattr(shape, 'height') else 0
            }

            # Classify by type
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                return SlideObject(
                    object_type="image",
                    description=f"Image: {shape.name}",
                    position=position,
                    size=size
                )
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                return SlideObject(
                    object_type="chart",
                    description=f"Chart: {shape.name}",
                    position=position,
                    size=size
                )
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                rows = len(shape.table.rows) if hasattr(shape, 'table') else 0
                cols = len(shape.table.columns) if hasattr(shape, 'table') else 0
                return SlideObject(
                    object_type="table",
                    description=f"Table: {rows}x{cols}",
                    position=position,
                    size=size
                )
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                # Count group as multiple objects
                return SlideObject(
                    object_type="group",
                    description=f"Group: {shape.name}",
                    position=position,
                    size=size
                )
            elif hasattr(shape, "text") and shape.text.strip():
                return SlideObject(
                    object_type="text",
                    description=f"Text: {shape.text[:50]}...",
                    position=position,
                    size=size
                )
            elif shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                return SlideObject(
                    object_type="shape",
                    description=f"Shape: {shape.name}",
                    position=position,
                    size=size
                )

            return None

        except Exception as e:
            logger.warning(f"Error classifying shape: {e}")
            return None

    def _calculate_complexity(
        self,
        object_count: int,
        has_image: bool,
        has_chart: bool,
        has_table: bool
    ) -> float:
        """Calculate slide complexity score"""
        score = object_count

        # Add weights for special elements
        if has_image:
            score += 2
        if has_chart:
            score += 3
        if has_table:
            score += 2

        return score

    def _generate_enhanced_markdown(
        self,
        prs: Presentation,
        pptx_path: Path,
        analyses: List[SlideAnalysis]
    ) -> str:
        """Generate enhanced markdown with metadata and images"""
        lines = []

        # Header
        lines.append(f"# {pptx_path.stem}\n")
        lines.append(f"*Converted from PowerPoint on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n")
        lines.append(f"**Total Slides:** {len(prs.slides)}\n")

        # Table of contents
        lines.append("## Table of Contents\n")
        for analysis in analyses:
            title = analysis.title or f"Slide {analysis.slide_number}"
            lines.append(f"- [Slide {analysis.slide_number}: {title}](#slide-{analysis.slide_number})\n")
        lines.append("\n---\n")

        # Process each slide
        for idx, (slide, analysis) in enumerate(zip(prs.slides, analyses), 1):
            lines.append(f"\n## Slide {idx}\n")

            if analysis.title:
                lines.append(f"### {analysis.title}\n")

            # Add metadata
            lines.append(f"**Objects:** {analysis.object_count} | "
                        f"**Complexity:** {analysis.complexity_score:.1f}\n")

            if analysis.should_save_as_image:
                # Save slide as image and reference it
                image_path = self._save_slide_as_image(slide, idx, pptx_path.stem)
                rel_path = Path("slides") / image_path.name
                lines.append(f"\n![Slide {idx}]({rel_path})\n")

                # Attempt Mermaid conversion if enabled
                if self.enable_mermaid:
                    mermaid = self._try_generate_mermaid(analysis)
                    if mermaid:
                        lines.append("\n### Diagram Representation\n")
                        lines.append("```mermaid\n")
                        lines.append(mermaid)
                        lines.append("\n```\n")
            else:
                # Extract text content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if shape == slide.shapes.title:
                            continue  # Already added as title
                        lines.append(f"\n{text}\n")

            # Add object inventory
            if analysis.object_count > 0:
                lines.append("\n**Slide Elements:**\n")
                for obj in analysis.objects[:10]:  # Limit to first 10
                    lines.append(f"- {obj.object_type}: {obj.description}\n")
                if len(analysis.objects) > 10:
                    lines.append(f"- ... and {len(analysis.objects) - 10} more\n")

            lines.append("\n---\n")

        return "".join(lines)

    def _save_slide_as_image(
        self,
        slide,
        slide_num: int,
        presentation_name: str
    ) -> Path:
        """Save a slide as PNG image"""
        try:
            # Note: python-pptx doesn't directly support slide-to-image conversion
            # This is a placeholder that creates a simple representation
            # In production, you'd use external tools like LibreOffice or PowerPoint API

            filename = f"{presentation_name}_slide_{slide_num:03d}.png"
            image_path = self.slides_dir / filename

            # Create a placeholder image with slide info
            img = Image.new('RGB', (1024, 768), color='white')
            draw = ImageDraw.Draw(img)

            # Add text indicating this is a placeholder
            text = f"Slide {slide_num}\n(Complex slide with {len(slide.shapes)} objects)\n\n"
            text += "Note: Actual slide rendering requires\nPowerPoint/LibreOffice integration"

            draw.text((50, 50), text, fill='black')
            img.save(image_path, 'PNG')

            logger.info(f"Saved slide {slide_num} as {image_path}")
            return image_path

        except Exception as e:
            logger.error(f"Error saving slide {slide_num} as image: {e}")
            raise

    def _try_generate_mermaid(self, analysis: SlideAnalysis) -> Optional[str]:
        """Attempt to generate Mermaid diagram from slide analysis"""
        if not self.enable_mermaid:
            return None

        # Heuristic-based Mermaid generation
        # In production, this would use LLM analysis

        if analysis.has_chart:
            # Generate simple chart representation
            return self._generate_chart_mermaid(analysis)

        # Check if slide looks like a flowchart
        shape_count = sum(1 for obj in analysis.objects if obj.object_type == "shape")
        if shape_count >= 3:
            return self._generate_flowchart_mermaid(analysis)

        return None

    def _generate_chart_mermaid(self, analysis: SlideAnalysis) -> str:
        """Generate Mermaid diagram for chart-like slides"""
        mermaid = "graph LR\n"
        mermaid += f"    A[{analysis.title or 'Chart Data'}]\n"

        for i, obj in enumerate(analysis.objects[:5], 1):
            if obj.object_type in ["chart", "shape", "text"]:
                label = obj.description.replace('"', "'")[:30]
                mermaid += f"    A --> B{i}[{label}]\n"

        return mermaid

    def _generate_flowchart_mermaid(self, analysis: SlideAnalysis) -> str:
        """Generate Mermaid flowchart for process-like slides"""
        mermaid = "flowchart TD\n"

        shapes = [obj for obj in analysis.objects if obj.object_type == "shape"]
        for i, shape in enumerate(shapes[:6], 1):
            label = shape.description.replace('"', "'")[:30]
            mermaid += f"    {chr(64+i)}[{label}]\n"

        # Add some connections
        for i in range(1, min(len(shapes), 5)):
            mermaid += f"    {chr(64+i)} --> {chr(65+i)}\n"

        return mermaid


class ObjectAnalyzer:
    """Specialized class for object analysis using LLM"""

    def __init__(self, llm_client, model: str = "gpt-4o"):
        self.llm_client = llm_client
        self.model = model

    def analyze_slide_image(self, image_path: Path) -> Dict[str, Any]:
        """
        Analyze a slide image using LLM for detailed object counting.

        Args:
            image_path: Path to slide image

        Returns:
            Dictionary with analysis results
        """
        try:
            # This would use the LLM client to analyze the image
            # Placeholder implementation
            return {
                "object_count": 0,
                "objects": [],
                "description": "Analysis not available",
                "mermaid": None
            }
        except Exception as e:
            logger.error(f"Error analyzing image {image_path}: {e}")
            return {
                "object_count": 0,
                "objects": [],
                "description": f"Error: {str(e)}",
                "mermaid": None
            }


def main():
    """Main entry point for testing"""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python enhanced_markdown_maker.py <pptx_file>")
        sys.exit(1)

    pptx_file = sys.argv[1]

    # Initialize without LLM client for basic testing
    maker = EnhancedMarkdownMaker(
        llm_client=None,
        output_dir="output",
        object_threshold=5
    )

    try:
        markdown, metadata = maker.convert_powerpoint(pptx_file)
        print(f"✅ Conversion complete!")
        print(f"📊 Total slides: {metadata.total_slides}")
        print(f"🖼️  Slides as images: {metadata.slides_as_images}")
        print(f"📝 Slides as text: {metadata.slides_as_text}")
        print(f"⏱️  Conversion time: {metadata.conversion_time:.2f}s")
    except Exception as e:
        print(f"❌ Error: {e}")
        logger.exception("Conversion failed")
        sys.exit(1)


if __name__ == "__main__":
    main()
